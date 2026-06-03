"""Genera la presentazione PowerPoint di Sertor (esperimento -> stato -> piano).

Uso: .venv-core/Scripts/python.exe presentation/build_pptx.py
Richiede `python-pptx` (installato nel venv, non è una dipendenza del progetto).
Output: presentation/sertor-overview.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

OUT = Path(__file__).parent / "sertor-overview.pptx"

# Palette
DARK = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x2D, 0x6C, 0xDF)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)
GREY = RGBColor(0x5F, 0x6B, 0x7A)
AMBER = RGBColor(0xB8, 0x6E, 0x00)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)

prs = Presentation()
prs.slide_width = Emu(12192000)   # 16:9
prs.slide_height = Emu(6858000)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _set(p, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Segoe UI"


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _band(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = color
    shp.shadow.inherit = False
    return shp


def title_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s, DARK)
    tf = _box(s, 900000, 2200000, 10400000, 1600000)
    _set(tf.paragraphs[0], "Sertor", 60, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    p = tf.add_paragraph()
    _set(p, "Framework RAG enterprise — dal prototipo alla produzione", 26, RGBColor(0xBF, 0xD2, 0xF5))
    tf2 = _box(s, 900000, 5200000, 10400000, 700000)
    _set(tf2.paragraphs[0],
         "Repo-agnostico · local-first ↔ cloud · libreria + CLI + MCP · governance SpecKit", 16,
         RGBColor(0x9A, 0xB0, 0xD0))


def section_slide(n, title, subtitle):
    s = prs.slides.add_slide(BLANK)
    _bg(s, LIGHT)
    _band(s, 700000, 2400000, 240000, 1600000, BLUE)
    tf = _box(s, 1100000, 2400000, 10000000, 1700000)
    _set(tf.paragraphs[0], f"Atto {n}", 20, BLUE, bold=True)
    p = tf.add_paragraph()
    _set(p, title, 40, DARK, bold=True)
    p2 = tf.add_paragraph()
    _set(p2, subtitle, 20, GREY)


def bullets_slide(title, bullets, foot=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, RGBColor(0xFF, 0xFF, 0xFF))
    _band(s, 0, 0, SW, 950000, DARK)
    tf = _box(s, 700000, 230000, 10800000, 600000)
    _set(tf.paragraphs[0], title, 30, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    body = _box(s, 800000, 1250000, 10600000, 4900000)
    first = True
    for level, text, color in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        marker = "•  " if level == 0 else "    –  "
        _set(p, marker + text, 22 if level == 0 else 18, color)
        p.space_after = Pt(8)
    if foot:
        ft = _box(s, 800000, 6200000, 10600000, 450000)
        _set(ft.paragraphs[0], foot, 14, GREY)


def kpi_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s, RGBColor(0xFF, 0xFF, 0xFF))
    _band(s, 0, 0, SW, 950000, DARK)
    tf = _box(s, 700000, 230000, 10800000, 600000)
    _set(tf.paragraphs[0], "Stato attuale — in numeri", 30, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    kpis = [("3 / 3", "Must del core (MVP) completi", GREEN),
            ("108", "test verdi (+2 xfail)", BLUE),
            ("9", "PR mergiate (pipeline SpecKit)", BLUE),
            ("147 / 1226", "documenti / chunk indicizzati", AMBER),
            ("9 / 9", "principi di Costituzione rispettati", GREEN),
            ("4", "bug trovati dal dogfooding", AMBER)]
    x0, y0, w, h, gx, gy = 800000, 1500000, 3500000, 1800000, 300000, 350000
    for i, (num, label, color) in enumerate(kpis):
        col, row = i % 3, i // 3
        card = _band(s, x0 + col * (w + gx), y0 + row * (h + gy), w, h, LIGHT)
        card.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEC)
        tfx = card.text_frame
        tfx.word_wrap = True
        _set(tfx.paragraphs[0], num, 40, color, bold=True, align=PP_ALIGN.CENTER)
        p = tfx.add_paragraph()
        _set(p, label, 15, GREY, align=PP_ALIGN.CENTER)


def arch_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s, RGBColor(0xFF, 0xFF, 0xFF))
    _band(s, 0, 0, SW, 950000, DARK)
    _set(_box(s, 700000, 230000, 10800000, 600000).paragraphs[0],
         "Architettura — Clean Architecture (dipendenze verso l'interno)", 26,
         RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    layers = [
        ("Consumatori:  CLI `sertor`  ·  Server MCP  ·  (futuro) altri client", BLUE, 1500000, 900000),
        ("Composition root  —  wiring da configurazione (.env)", GREY, 2550000, 750000),
        ("Adapters:  Ollama / Azure (embeddings, LLM)  ·  Chroma / Azure AI Search (store)",
         AMBER, 3450000, 800000),
        ("Services / Engines / Wiki:  ingestione · chunking · indexing · retrieval · baseline · wiki",
         GREEN, 4400000, 800000),
        ("Domain (cuore):  entità · porte (EmbeddingProvider, VectorStore, LLMProvider) · errori",
         DARK, 5350000, 850000),
    ]
    for text, color, y, h in layers:
        band = _band(s, 900000, y, 10400000, h, LIGHT)
        band.line.color.rgb = color
        _set(band.text_frame.paragraphs[0], text, 16, color, bold=True, align=PP_ALIGN.CENTER)
    _set(_box(s, 900000, 6250000, 10400000, 400000).paragraphs[0],
         "Il core non importa SDK di provider né la CLI; gli adapter dipendono dalle astrazioni; "
         "il wiring sta nel composition root.", 13, GREY)


def roadmap_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s, RGBColor(0xFF, 0xFF, 0xFF))
    _band(s, 0, 0, SW, 950000, DARK)
    _set(_box(s, 700000, 230000, 10800000, 600000).paragraphs[0],
         "Piano futuro — roadmap", 30, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    cols = [
        ("✓ Fatto (in produzione)", GREEN,
         ["Nucleo di retrieval (FEAT-001)", "Motore baseline (FEAT-002)", "Skill LLM Wiki (FEAT-003)",
          "CLI `sertor` (esecuzione)", "Server MCP sul motore nuovo", "Dogfooding su Azure"]),
        ("▶ Prossimo (Should)", BLUE,
         ["RAG ibrido + reranking (FEAT-004)", "RAG a grafo / GraphRAG (FEAT-005)",
          "RAG agentico (FEAT-006)", "Manutenzione wiki (FEAT-007)", "CLI: install selettivo"]),
        ("◌ Dopo (Could / Won't ora)", GREY,
         ["Arricchimento Wiki↔RAG (FEAT-008)", "Refresh incrementale (FEAT-009)",
          "CLI: wizard config", "Distribuzione PyPI"]),
    ]
    x0, w, gx = 700000, 3600000, 250000
    for i, (head, color, items) in enumerate(cols):
        x = x0 + i * (w + gx)
        hb = _band(s, x, 1300000, w, 650000, color)
        _set(hb.text_frame.paragraphs[0], head, 17, RGBColor(0xFF, 0xFF, 0xFF), bold=True,
             align=PP_ALIGN.CENTER)
        card = _band(s, x, 2050000, w, 3900000, LIGHT)
        card.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEC)
        tf = card.text_frame
        tf.word_wrap = True
        first = True
        for it in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            _set(p, "•  " + it, 15, DARK)
            p.space_after = Pt(7)


# --- costruzione del deck ---
title_slide()
bullets_slide("Cos'è Sertor", [
    (0, "In una frase: un motore di ricerca semantica per il tuo progetto — codice + documentazione — portabile e governato.", DARK),
    (0, "Indicizza un repository e recupera i passaggi più pertinenti a una domanda in linguaggio naturale, citando il file.", GREY),
    (1, "RAG = Retrieval-Augmented Generation; Sertor copre la parte di retrieval (trovare il contesto giusto).", GREY),
    (0, "Tre capacità:", DARK),
    (1, "Nucleo di retrieval riusabile: ingestione, chunking, embeddings, vector store, facade.", GREY),
    (1, "Motori RAG: baseline oggi; ibrido / grafo / agentico in arrivo.", GREY),
    (1, "Skill LLM Wiki: la conoscenza distillata (decisioni, il “perché”) indicizzata col resto.", GREY),
    (0, "Come si usa: libreria Python · CLI `sertor` · strumento per agenti LLM (via MCP).", DARK),
    (0, "Cosa NON è: non è un chatbot né un LLM proprio; non produce la risposta finale — fornisce il contesto pertinente.", AMBER),
], foot="Obiettivo: portare questa capacità su QUALUNQUE repo, in modo riproducibile e production-grade.")
bullets_slide("Il problema & la visione", [
    (0, "La conoscenza di un progetto (codice + perché/decisioni) si disperde e si ricostruisce ogni volta.", GREY),
    (0, "Obiettivo: portare capacità RAG su QUALUNQUE repository, in modo riproducibile.", DARK),
    (1, "Repo-agnostico, production-grade, testabile e misurabile.", GREY),
    (1, "Local-first ↔ cloud commutabile via configurazione (no lock-in).", GREY),
    (1, "Riusabile come libreria, esposto via CLI e via MCP (per agenti LLM).", GREY),
], foot="Sertor = nucleo di retrieval + motori RAG + skill LLM Wiki, sopra una governance esplicita.")

section_slide(1, "L'esperimento", "Il prototipo: 4 approcci RAG su un corpus reale")
bullets_slide("Atto 1 — Il prototipo (congelato)", [
    (0, "4 modalità RAG costruite su un corpus FastAPI:", DARK),
    (1, "Baseline vettoriale · Ibrido (BM25+dense+reranking) · GraphRAG · Agentico.", GREY),
    (0, "Local-first (Ollama + Chroma) e variante Azure (OpenAI + AI Search).", GREY),
    (0, "Ha dimostrato la fattibilità tecnica delle 4 strade.", GREEN),
    (0, "Ma: esplorativo, codice accoppiato al corpus, senza test, non riusabile come componente.", AMBER),
    (1, "Lezione → serve rifarlo a qualità di produzione, repo-agnostico e testato.", DARK),
], foot="Il prototipo è ora congelato e indicizzato come riferimento (dogfooding).")

section_slide(2, "La svolta in produzione", "Governance esplicita + Clean Architecture")
bullets_slide("Atto 2 — Come lavoriamo ora", [
    (0, "Costituzione di progetto (v1.0.0): 9 principi vincolanti.", DARK),
    (1, "Es. core a dipendenze verso l'interno; errori espliciti; idempotenza; config centralizzata; osservabilità.", GREY),
    (0, "Pipeline SpecKit per ogni feature, con gate di qualità:", DARK),
    (1, "requirements (EARS) → spec → plan (Constitution Check) → tasks → analyze → implement.", GREY),
    (0, "Clean Architecture: dominio puro ← servizi/adapter; provider intercambiabili dietro porte.", GREY),
    (0, "Branch + PR + test per ogni passo (niente push diretti su master).", GREEN),
])
arch_slide()

section_slide(3, "Stato attuale", "MVP del core in produzione, dogfoodato su sé stesso")
bullets_slide("Atto 3 — Cosa funziona oggi", [
    (0, "MVP del core completo e in `master`:", GREEN),
    (1, "Nucleo di retrieval · Motore RAG baseline · Skill LLM Wiki (crea + indicizza).", GREY),
    (0, "CLI `sertor`: `index` / `search` / `wiki index` + osservabilità (log/Splunk).", DARK),
    (0, "Server MCP sul motore nuovo: Claude interroga Sertor come strumento.", DARK),
    (0, "Dogfooding di produzione: indicizzato Sertor stesso su Azure (embeddings) + Chroma.", BLUE),
    (1, "Il dogfooding reale ha scovato 4 bug invisibili ai test con mock → tutti corretti con regressioni.", AMBER),
])
kpi_slide()

section_slide(4, "Piano futuro", "Dai motori avanzati alla distribuzione")
roadmap_slide()
bullets_slide("Questioni aperte & come contribuire", [
    (0, "Aperte (per scelta):", DARK),
    (1, "Fissare le soglie di pertinenza misurando su un ground-truth reale (ora possibile su Azure).", GREY),
    (1, "Promuovere PowerShell / SQL da chunking di fallback a sintattico.", GREY),
    (0, "Come usarlo (per il team):", DARK),
    (1, "`uv pip install -e .` → comando `sertor`; oppure i tool MCP `search_code/docs/combined` in Claude.", GREY),
    (1, "Ogni nuova capacità segue la stessa pipeline SpecKit, sotto Costituzione.", GREEN),
], foot="Domande? — il wiki di progetto documenta design, implementazione e dogfooding in modo cumulativo.")

prs.save(str(OUT))
print(f"OK: {OUT}  ({len(prs.slides)} slide)")
